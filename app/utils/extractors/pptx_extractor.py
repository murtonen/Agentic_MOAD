import os
from pptx import Presentation
from typing import Dict, Any, List
import re

class PPTXExtractor:
    """Class to extract text content from PowerPoint files with improved structure preservation."""
    
    def __init__(self, file_path: str):
        """
        Initialize with the path to the PowerPoint file.
        
        Args:
            file_path: Path to the PowerPoint file
        """
        self.file_path = file_path
        
    def extract_all_slides(self) -> Dict[str, str]:
        """
        Extract text content from all slides in the PowerPoint file.
        
        Returns:
            Dictionary mapping slide identifiers to slide content
        """
        return self.extract_content()
        
    def extract_content(self) -> Dict[str, str]:
        """
        Extract text content from all slides in the PowerPoint file.
        
        Returns:
            Dictionary mapping slide identifiers to slide content
        """
        if not os.path.exists(self.file_path):
            raise FileNotFoundError(f"PowerPoint file not found: {self.file_path}")
        
        try:
            presentation = Presentation(self.file_path)
        except Exception as e:
            raise ValueError(f"Error opening PowerPoint file: {str(e)}")
        
        content = {}
        
        for i, slide in enumerate(presentation.slides):
            slide_id = f"slide_{i+1}"
            try:
                slide_text = self._extract_slide_text_with_structure(slide)
                
                # Only include slides with actual content
                if slide_text.strip():
                    content[slide_id] = slide_text
            except Exception as e:
                print(f"Error extracting content from slide {i+1}: {str(e)}")
                # Add a simple version for slides that cause errors
                content[slide_id] = f"[Error extracting content from slide {i+1}]"
        
        return content
    
    def _extract_slide_text(self, slide) -> str:
        """
        Extract text from a single slide (basic method).
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            Extracted text content
        """
        text_parts = []
        
        # Extract title if present
        if slide.shapes.title and slide.shapes.title.text:
            text_parts.append(f"Title: {slide.shapes.title.text}")
        
        # Extract text from all shapes
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
                
            if shape.text and shape != slide.shapes.title:
                text_parts.append(shape.text)
        
        return "\n".join(text_parts)
    
    def _extract_slide_text_with_structure(self, slide) -> str:
        """
        Extract text with improved structure preservation from a single slide.
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            Extracted text content with structure information
        """
        text_parts = []
        slide_title = ""
        
        # Extract title if present (store separately)
        if slide.shapes.title and slide.shapes.title.text:
            slide_title = slide.shapes.title.text.strip()
            text_parts.append(f"Title: {slide_title}")
        
        # Look for tables and extract them with structure
        tables = []
        bullets = []
        plain_text = []
        
        for shape in slide.shapes:
            # Skip the title shape that we already processed
            if slide.shapes.title and shape == slide.shapes.title:
                continue
            
            # Check if it's a table using the shape type
            try:
                # In python-pptx 1.0.2, we should check for shape type or if it has a table property
                if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                    table_data = self._extract_table(shape.table)
                    if table_data:
                        tables.append(table_data)
                    continue
            except (AttributeError, ValueError):
                # If we can't access the table property or it raises an error, continue to next shape
                pass
                
            # Extract text from shape
            if hasattr(shape, "text") and shape.text.strip():
                shape_text = shape.text.strip()
                
                # Check if text appears to be a bulleted list
                if self._is_bullet_list(shape_text):
                    bullets.append(self._format_bullets(shape_text))
                else:
                    plain_text.append(shape_text)
        
        # Add tables with structure
        if tables:
            text_parts.append("\n--- Tables ---")
            for i, table in enumerate(tables):
                text_parts.append(f"Table {i+1}:")
                text_parts.append(table)
        
        # Add bulleted lists
        if bullets:
            text_parts.append("\n--- Lists ---")
            for bullet_list in bullets:
                text_parts.append(bullet_list)
        
        # Add remaining text
        if plain_text:
            text_parts.append("\n--- Content ---")
            for text in plain_text:
                text_parts.append(text)
        
        # Add metadata based on slide title
        metadata = self._extract_metadata(slide_title, "".join(plain_text))
        if metadata:
            text_parts.append("\n--- Metadata ---")
            for key, value in metadata.items():
                text_parts.append(f"{key}: {value}")
        
        return "\n\n".join(text_parts)
    
    def _extract_table(self, table) -> str:
        """
        Extract data from a table with structure.
        
        Args:
            table: PowerPoint table object
            
        Returns:
            Formatted table as string
        """
        try:
            rows = []
            for row in table.rows:
                cells = []
                for cell in row.cells:
                    if cell.text.strip():
                        cells.append(cell.text.strip())
                    else:
                        cells.append("-")
                rows.append(" | ".join(cells))
            
            return "\n".join(rows)
        except Exception as e:
            print(f"Error extracting table data: {str(e)}")
            return "[Error extracting table data]"
    
    def _is_bullet_list(self, text: str) -> bool:
        """
        Check if text appears to be a bullet list.
        
        Args:
            text: The text to check
            
        Returns:
            True if it appears to be a bullet list
        """
        lines = text.split('\n')
        if len(lines) < 2:
            return False
            
        # Check for bullet patterns
        bullet_patterns = [r'^\s*[\•\-\*]\s', r'^\s*\d+[\.\)]\s']
        bullet_count = 0
        
        for line in lines:
            if any(re.match(pattern, line) for pattern in bullet_patterns):
                bullet_count += 1
                
        # If more than half the lines match a bullet pattern, consider it a list
        return bullet_count >= len(lines) / 2
    
    def _format_bullets(self, text: str) -> str:
        """
        Format bullet list with proper structure.
        
        Args:
            text: The bullet list text
            
        Returns:
            Formatted bullet list
        """
        return text  # Already has structure from PowerPoint
    
    def _extract_metadata(self, title: str, content: str) -> Dict[str, str]:
        """
        Extract metadata from slide title and content.
        
        Args:
            title: Slide title
            content: Slide content
            
        Returns:
            Dictionary of metadata
        """
        metadata = {}
        
        # Check for license or product information
        license_patterns = [
            (r'\b(standard|pro|enterprise|pro\+)\b', 'license_tier'),
            (r'\b(itsm|csx|hrsd|itom)\b', 'product_line'),
            (r'\b(virtual agent|now assist|ai|workflow)\b', 'feature_category')
        ]
        
        combined_text = (title + " " + content).lower()
        
        for pattern, meta_key in license_patterns:
            matches = re.findall(pattern, combined_text)
            if matches:
                metadata[meta_key] = ", ".join(set(matches))
        
        # Detect if it's a comparison chart
        if ('compare' in combined_text or 'vs' in combined_text or 'versus' in combined_text) and any(term in combined_text for term in ['edition', 'license', 'tier']):
            metadata['content_type'] = 'comparison_chart'
        
        # Detect capability matrices
        if 'capability' in combined_text and 'matrix' in combined_text:
            metadata['content_type'] = 'capability_matrix'
        
        return metadata 